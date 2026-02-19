"""
공통 슬라이드 유틸리티
레퍼런스 PPT의 공통 요소를 함수화
"""
import os
from typing import Optional

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def add_group_marker(slide):
    """
    둥근사각 2개 겹친 Group 장식 마커 (70% 축소)
    원본 크기 약 0.312"×0.299" → 70% 약 0.218"×0.209"
    """
    # 70% 스케일
    # Rect1: roundRect, accent3, 원본 (163375, 253000, 200025, 190500)
    rect1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(114363), Emu(177100),
        Emu(140018), Emu(133350),
    )
    rect1.fill.solid()
    rect1.fill.fore_color.rgb = RGBColor(0xEE, 0xFF, 0x41)
    rect1.line.fill.background()

    # Rect2: roundRect, #221E1F, 원본 (315325, 108450, 200025, 190500)
    rect2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(220728), Emu(75915),
        Emu(140018), Emu(133350),
    )
    rect2.fill.solid()
    rect2.fill.fore_color.rgb = RGBColor(0x22, 0x1E, 0x1F)
    rect2.line.fill.background()


def add_slide_header(slide, title: str, subtitle: str):
    """
    공통 슬라이드 헤더
    - Group 마커 + 21pt bold 제목 (0.428", 0.123")
    - 16pt bold 서브타이틀 (4.371", 0.206", RIGHT 정렬)
    - 구분선 없음
    """
    # Group 장식 마커
    add_group_marker(slide)

    # 제목 텍스트 (21pt bold)
    txBox = slide.shapes.add_textbox(
        Inches(0.428), Inches(0.123),
        Inches(3.5), Inches(0.45),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x22, 0x1E, 0x1F)

    # 서브타이틀 (16pt bold, 우측정렬)
    txBox2 = slide.shapes.add_textbox(
        Inches(4.371), Inches(0.206),
        Inches(5.2), Inches(0.4),
    )
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0x22, 0x1E, 0x1F)
    p2.alignment = PP_ALIGN.RIGHT


def add_logo(slide, logo_path: Optional[str] = None):
    """
    중앙 하단 로고 이미지
    위치: (4.139", 5.726", 1.678"×0.478")
    """
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(
            logo_path,
            Inches(4.139), Inches(5.726),
            Inches(1.678), Inches(0.478),
        )


def add_source_text(slide, text: str):
    """
    출처 텍스트 (좌하단)
    위치: (0.328", 5.686"), 10pt
    """
    txBox = slide.shapes.add_textbox(
        Inches(0.328), Inches(5.686),
        Inches(6), Inches(0.3),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
