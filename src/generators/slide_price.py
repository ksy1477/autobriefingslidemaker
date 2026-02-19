"""
실거래가 슬라이드 생성
"""
import os
from typing import Optional, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from src.models import PriceInfo


def add_price_slide(
    prs: Presentation,
    complex_name: str,
    price_info: PriceInfo,
    price_summary_text: str,
    logo_path: Optional[str] = None,
):
    """
    실거래가 슬라이드 추가

    Args:
        prs: Presentation 객체
        complex_name: 단지명
        price_info: 실거래가 정보
        price_summary_text: 요약 텍스트
        logo_path: 회사 로고 경로
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 헤더
    _add_slide_header(slide, complex_name, "실거래가")

    # 요약 텍스트
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(9), Inches(1.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(price_summary_text.split("\n")):
        if i == 0:
            tf.paragraphs[0].text = line
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            tf.paragraphs[0].space_after = Pt(4)
        else:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.space_after = Pt(4)

    # 최근 실거래 테이블 (좌측)
    transactions = price_info.recent_transactions[:6]  # 최대 6건
    if transactions:
        cols = 4  # 날짜, 평형, 층, 가격
        rows = len(transactions) + 1
        table_shape = slide.shapes.add_table(
            rows, cols,
            Inches(0.3), Inches(2.8),
            Inches(4.5), Inches(min(3.0, rows * 0.35))
        )
        table = table_shape.table

        col_widths = [Inches(1.2), Inches(0.9), Inches(0.7), Inches(1.7)]
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

        # 헤더
        headers = ["거래일", "평형", "층", "거래가"]
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xC8, 0x10, 0x2E)

        # 데이터
        for row_idx, txn in enumerate(transactions, start=1):
            data = [
                txn.date.strftime("%Y.%m.%d"),
                txn.area_pyeong,
                f"{txn.floor}층",
                txn.price,
            ]
            for col_idx, val in enumerate(data):
                cell = table.cell(row_idx, col_idx)
                cell.text = val
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(9)
                    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    p.alignment = PP_ALIGN.CENTER
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    # 매매가 추이 그래프 (우측 - 이미지)
    if price_info.price_graph_image_path and os.path.exists(price_info.price_graph_image_path):
        slide.shapes.add_picture(
            price_info.price_graph_image_path,
            Inches(5.0), Inches(2.8), Inches(4.7), Inches(3.0)
        )

    # 출처
    txBox2 = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5), Inches(5), Inches(0.3)
    )
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "*아실 (https://asil.kr)  *국토교통부 실거래가 공개시스템"
    p2.font.size = Pt(8)
    p2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # 로고
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(
            logo_path, Inches(7.5), Inches(6.2), Inches(2), Inches(0.6)
        )

    return slide


def _add_slide_header(slide, title: str, subtitle: str):
    """슬라이드 공통 헤더"""
    marker = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.35), Pt(8), Inches(0.4)
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = RGBColor(0xC8, 0x10, 0x2E)
    marker.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.3), Inches(6), Inches(0.5)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    txBox2 = slide.shapes.add_textbox(
        Inches(7.5), Inches(0.3), Inches(2), Inches(0.5)
    )
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p2.alignment = PP_ALIGN.RIGHT

    line = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.85), Inches(9.4), Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    line.line.fill.background()
