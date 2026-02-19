"""
입지정보 슬라이드 생성 (Slide B)
최근접 지하철역 도보 경로 + 강남역 대중교통 경로
"""
import os
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from src.models import LocationInfo


def add_location_slide(
    prs: Presentation,
    complex_name: str,
    location_info: LocationInfo,
    logo_path: Optional[str] = None,
):
    """
    입지정보 슬라이드 추가

    Args:
        prs: Presentation 객체
        complex_name: 단지명
        location_info: 입지 정보
        logo_path: 회사 로고 경로
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 헤더
    _add_slide_header(slide, complex_name, "입지정보")

    # ── 역 정보 텍스트 ──
    station_text = ""
    if location_info.nearest_station:
        station_text = f"{location_info.nearest_station}"
        if location_info.station_line:
            station_text += f"({location_info.station_line})"
        station_text += f" - 도보 {location_info.walk_minutes}분"

    gangnam_text = f"강남역 - 대중교통 {location_info.gangnam_minutes}분"

    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1), Inches(9), Inches(0.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    if station_text:
        p = tf.paragraphs[0]
        p.text = station_text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = gangnam_text
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    else:
        tf.paragraphs[0].text = gangnam_text
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── 도보 경로 지도 (좌측) ──
    walk_label = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.0), Inches(4.3), Inches(0.3)
    )
    walk_label.text_frame.paragraphs[0].text = "도보 경로"
    walk_label.text_frame.paragraphs[0].font.size = Pt(10)
    walk_label.text_frame.paragraphs[0].font.bold = True
    walk_label.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xC8, 0x10, 0x2E)

    if (location_info.walk_route_image_path
            and os.path.exists(location_info.walk_route_image_path)):
        slide.shapes.add_picture(
            location_info.walk_route_image_path,
            Inches(0.5), Inches(2.3), Inches(4.3), Inches(3.5)
        )
    else:
        _add_placeholder_text(slide, Inches(0.5), Inches(2.3), Inches(4.3), Inches(3.5),
                              "[도보 경로 지도]")

    # ── 대중교통 경로 지도 (우측) ──
    transit_label = slide.shapes.add_textbox(
        Inches(5.2), Inches(2.0), Inches(4.3), Inches(0.3)
    )
    transit_label.text_frame.paragraphs[0].text = "대중교통 경로 (→ 강남역)"
    transit_label.text_frame.paragraphs[0].font.size = Pt(10)
    transit_label.text_frame.paragraphs[0].font.bold = True
    transit_label.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xC8, 0x10, 0x2E)

    if (location_info.transit_route_image_path
            and os.path.exists(location_info.transit_route_image_path)):
        slide.shapes.add_picture(
            location_info.transit_route_image_path,
            Inches(5.2), Inches(2.3), Inches(4.3), Inches(3.5)
        )
    else:
        _add_placeholder_text(slide, Inches(5.2), Inches(2.3), Inches(4.3), Inches(3.5),
                              "[대중교통 경로 지도]")

    # ── 출처 ──
    txBox3 = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5), Inches(5), Inches(0.3)
    )
    p3 = txBox3.text_frame.paragraphs[0]
    p3.text = "*네이버지도 (https://map.naver.com)"
    p3.font.size = Pt(8)
    p3.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # ── 로고 ──
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(
            logo_path, Inches(7.5), Inches(6.2), Inches(2), Inches(0.6)
        )

    return slide


def _add_placeholder_text(slide, left, top, width, height, text):
    """이미지 없을 때 텍스트 placeholder"""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER


def _add_slide_header(slide, title: str, subtitle: str):
    """슬라이드 공통 헤더 (제목 + 서브타이틀)"""
    # 좌측 빨간 마커
    marker = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.35), Pt(8), Inches(0.4)
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = RGBColor(0xC8, 0x10, 0x2E)
    marker.line.fill.background()

    # 단지명
    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.3), Inches(6), Inches(0.5)
    )
    p = txBox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 우측 서브타이틀
    txBox2 = slide.shapes.add_textbox(
        Inches(7.5), Inches(0.3), Inches(2), Inches(0.5)
    )
    p2 = txBox2.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p2.alignment = PP_ALIGN.RIGHT

    # 구분선
    line = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.85), Inches(9.4), Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    line.line.fill.background()
