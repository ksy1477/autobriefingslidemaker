"""
단지 개요 슬라이드 생성
"""
import os
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from src.models import ComplexInfo


def add_overview_slide(
    prs: Presentation,
    complex_info: ComplexInfo,
    overview_text: str,
    logo_path: Optional[str] = None,
):
    """
    단지 개요 슬라이드 추가

    Args:
        prs: Presentation 객체
        complex_info: 단지 정보
        overview_text: 개요 텍스트 (data_aggregator에서 생성)
        logo_path: 회사 로고 경로
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 헤더 (단지명 + "개요")
    _add_slide_header(slide, complex_info.name, "개요")

    # 개요 텍스트
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(4.5), Inches(2.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(overview_text.split("\n")):
        if i == 0:
            tf.paragraphs[0].text = line
            tf.paragraphs[0].font.size = Pt(13)
            tf.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            tf.paragraphs[0].space_after = Pt(6)
        else:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.space_after = Pt(6)

    # 해시태그
    if complex_info.hashtags:
        hashtag_text = "  ".join([f"#{tag}" for tag in complex_info.hashtags])
        txBox2 = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.8), Inches(4.5), Inches(0.4)
        )
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = hashtag_text
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(0xC8, 0x10, 0x2E)

    # 전경사진 (있으면)
    if complex_info.aerial_photo_path and os.path.exists(complex_info.aerial_photo_path):
        slide.shapes.add_picture(
            complex_info.aerial_photo_path,
            Inches(5.2), Inches(1.3), Inches(4.3), Inches(2.5)
        )

    # 위성지도 (있으면)
    if complex_info.satellite_map_path and os.path.exists(complex_info.satellite_map_path):
        slide.shapes.add_picture(
            complex_info.satellite_map_path,
            Inches(5.2), Inches(4.0), Inches(4.3), Inches(2.5)
        )

    # 출처
    txBox3 = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5), Inches(5), Inches(0.3)
    )
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "*네이버부동산 (https://land.naver.com)  *네이버지도 (https://map.naver.com)"
    p3.font.size = Pt(8)
    p3.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # 로고
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(
            logo_path, Inches(7.5), Inches(6.2), Inches(2), Inches(0.6)
        )

    return slide


def _add_slide_header(slide, title: str, subtitle: str):
    """슬라이드 공통 헤더 (제목 + 서브타이틀)"""
    # 좌측 빨간 마커
    marker = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.35), Pt(8), Inches(0.4)
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = RGBColor(0xC8, 0x10, 0x2E)
    marker.line.fill.background()

    # 단지명
    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.3), Inches(6), Inches(0.5)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 우측 서브타이틀
    txBox2 = slide.shapes.add_textbox(
        Inches(7.5), Inches(0.3), Inches(2), Inches(0.5)
    )
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p2.alignment = PP_ALIGN.RIGHT

    # 구분선
    line = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.85), Inches(9.4), Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    line.line.fill.background()
