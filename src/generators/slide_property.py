"""
매물정보 슬라이드 생성 — 레퍼런스 PPT 양식
"""
import os
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from src.models import PropertyDetail
from src.generators.slide_utils import add_slide_header, add_logo, add_source_text


def add_property_slide(
    prs: Presentation,
    prop: PropertyDetail,
    logo_path: Optional[str] = None,
):
    """
    매물정보 슬라이드 추가 (레퍼런스 레이아웃)

    - 공통 헤더
    - 빨간 세로 마커 + "단지내 위치" 라벨 at (0.127", 1.255")
    - 배치도 이미지: (0.084", 1.688", 2.712"×3.673")
    - 평면도: (2.864", 1.688", 4.607"×3.673")
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 공통 헤더
    title = f"{prop.complex_name} - {prop.dong} {prop.floor}"
    add_slide_header(slide, title, "매물정보")

    # ── 빨간 세로 마커 + "단지내 위치" 라벨 ──
    _add_red_marker(slide, Inches(0.127), Inches(1.255))
    label_box = slide.shapes.add_textbox(
        Inches(0.22), Inches(1.225),
        Inches(2), Inches(0.3),
    )
    lp = label_box.text_frame.paragraphs[0]
    lp.text = "단지내 위치"
    lp.font.size = Pt(12)
    lp.font.bold = True
    lp.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── 빨간 세로 마커 + "평면도" 라벨 ──
    _add_red_marker(slide, Inches(2.864), Inches(1.255))
    label_box2 = slide.shapes.add_textbox(
        Inches(2.96), Inches(1.225),
        Inches(2), Inches(0.3),
    )
    lp2 = label_box2.text_frame.paragraphs[0]
    lp2.text = "평면도"
    lp2.font.size = Pt(12)
    lp2.font.bold = True
    lp2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── 좌측: 배치도/동 위치 이미지 (0.084", 1.688", 2.712"×3.673") ──
    if prop.dong_location_image_path and os.path.exists(prop.dong_location_image_path):
        slide.shapes.add_picture(
            prop.dong_location_image_path,
            Inches(0.084), Inches(1.688),
            Inches(2.712), Inches(3.673),
        )
    else:
        _add_placeholder_text(
            slide,
            Inches(0.084), Inches(1.688),
            Inches(2.712), Inches(3.673),
            "[단지 내 위치]",
        )

    # ── 우측: 평면도 (2.864", 1.688", 4.607"×3.673") ──
    if prop.floor_plan_image_path and os.path.exists(prop.floor_plan_image_path):
        slide.shapes.add_picture(
            prop.floor_plan_image_path,
            Inches(2.864), Inches(1.688),
            Inches(4.607), Inches(3.673),
        )
    else:
        _add_placeholder_text(
            slide,
            Inches(2.864), Inches(1.688),
            Inches(4.607), Inches(3.673),
            "[평면도]",
        )

    # ── 빨간 세로 마커 + "매물 상세" 라벨 ──
    _add_red_marker(slide, Inches(7.54), Inches(1.255))
    label_box3 = slide.shapes.add_textbox(
        Inches(7.635), Inches(1.225),
        Inches(2), Inches(0.3),
    )
    lp3 = label_box3.text_frame.paragraphs[0]
    lp3.text = "매물 상세"
    lp3.font.size = Pt(12)
    lp3.font.bold = True
    lp3.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── 매물 상세 정보 테이블 (우측) ──
    info_rows = []
    info_rows.append(("매매가", prop.price))
    info_rows.append(("동 / 층", f"{prop.dong} / {prop.floor}"))
    if prop.area_m2:
        info_rows.append(("전용면적", f"{prop.area_m2}㎡"))
    if prop.direction:
        info_rows.append(("향", prop.direction))
    if prop.structure:
        info_rows.append(("구조", prop.structure))
    if prop.rooms is not None:
        bath_text = f" / 화장실 {prop.bathrooms}개" if prop.bathrooms else ""
        info_rows.append(("방/화장실", f"방 {prop.rooms}개{bath_text}"))
    if prop.memo:
        info_rows.append(("특이사항", prop.memo))

    row_count = len(info_rows)
    tbl_left = Inches(7.54)
    tbl_top = Inches(1.688)
    tbl_width = Inches(2.3)
    row_height = Inches(0.38)
    tbl_height = row_height * row_count

    table_shape = slide.shapes.add_table(
        row_count, 2,
        tbl_left, tbl_top,
        tbl_width, tbl_height,
    )
    table = table_shape.table
    table.columns[0].width = Inches(0.78)
    table.columns[1].width = Inches(1.52)

    BRAND_RED = RGBColor(0xC8, 0x10, 0x2E)
    LABEL_BG = RGBColor(0xF5, 0xF5, 0xF5)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x33, 0x33, 0x33)
    GRAY = RGBColor(0x55, 0x55, 0x55)
    BORDER = RGBColor(0xE0, 0xE0, 0xE0)

    for row_idx, (label, value) in enumerate(info_rows):
        # 라벨 셀 (좌)
        label_cell = table.cell(row_idx, 0)
        label_cell.text = label
        label_cell.fill.solid()
        label_cell.fill.fore_color.rgb = LABEL_BG
        for p in label_cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = DARK
            p.alignment = PP_ALIGN.CENTER
        label_cell.text_frame.word_wrap = True
        label_cell.margin_left = Pt(4)
        label_cell.margin_right = Pt(4)
        label_cell.margin_top = Pt(2)
        label_cell.margin_bottom = Pt(2)

        # 값 셀 (우)
        value_cell = table.cell(row_idx, 1)
        value_cell.text = value
        value_cell.fill.solid()
        value_cell.fill.fore_color.rgb = WHITE
        is_price = (label == "매매가")
        for p in value_cell.text_frame.paragraphs:
            p.font.size = Pt(9) if is_price else Pt(8)
            p.font.bold = is_price
            p.font.color.rgb = BRAND_RED if is_price else GRAY
        value_cell.text_frame.word_wrap = True
        value_cell.margin_left = Pt(6)
        value_cell.margin_right = Pt(4)
        value_cell.margin_top = Pt(2)
        value_cell.margin_bottom = Pt(2)

    # 로고
    add_logo(slide, logo_path)

    return slide


def _add_red_marker(slide, left, top):
    """빨간 세로 마커 (0.052"×0.197")"""
    marker = slide.shapes.add_shape(
        1,  # RECTANGLE
        left, top,
        Inches(0.052), Inches(0.197),
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = RGBColor(0xC8, 0x10, 0x2E)
    marker.line.fill.background()


def _add_placeholder_text(slide, left, top, width, height, text):
    """이미지 없을 때 텍스트 placeholder"""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER
