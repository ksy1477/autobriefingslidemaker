"""
학군지도 슬라이드 생성 — 레퍼런스 PPT 양식
초등학교 학구도 (Slide C) + 중·고등학교 학군 지도 (Slide D)
"""
import os
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from src.models import SchoolInfo
from src.generators.slide_utils import add_slide_header, add_logo, add_source_text


def add_elementary_school_slide(
    prs: Presentation,
    complex_name: str,
    school_info: SchoolInfo,
    logo_path: Optional[str] = None,
):
    """
    학군지도(초등학교) 슬라이드 추가 (Slide C)

    - 공통 헤더 (Group 마커, 구분선 없음)
    - 왼쪽: 학구도 지도 이미지
    - 오른쪽: 배정 초등학교 상세정보 테이블
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 공통 헤더
    add_slide_header(slide, "학구도(초등학교)", "학군지도")

    # ── 학구도 지도 이미지 (왼쪽) ──
    if (school_info.elementary_map_path
            and os.path.exists(school_info.elementary_map_path)):
        slide.shapes.add_picture(
            school_info.elementary_map_path,
            Inches(0.5), Inches(1.2), Inches(5.0), Inches(4.5)
        )
    else:
        _add_placeholder_text(
            slide, Inches(0.5), Inches(1.2), Inches(5.0), Inches(4.5),
            "[초등학교 학구도 지도]"
        )

    # ── 배정 초등학교 상세정보 (오른쪽) ──
    _add_school_detail_table(slide, school_info)

    # 출처
    add_source_text(slide, "*네이버지도  *NEIS 학교정보")

    # 로고
    add_logo(slide, logo_path)

    return slide


def _add_school_detail_table(slide, school_info: SchoolInfo):
    """배정 초등학교 상세정보 테이블 추가 (오른쪽 영역)"""

    # 제목: "배정 초등학교"
    title_box = slide.shapes.add_textbox(
        Inches(5.8), Inches(1.2), Inches(3.8), Inches(0.4)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "배정 초등학교"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xC8, 0x10, 0x2E)

    # 학교명
    if school_info.elementary_name:
        name_box = slide.shapes.add_textbox(
            Inches(5.8), Inches(1.55), Inches(3.8), Inches(0.35)
        )
        tf2 = name_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = school_info.elementary_name
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 도보거리
    if school_info.elementary_walk_distance:
        tab_box = slide.shapes.add_textbox(
            Inches(5.8), Inches(1.95), Inches(3.8), Inches(0.4)
        )
        tf_tab = tab_box.text_frame
        tf_tab.word_wrap = True
        p_tab = tf_tab.paragraphs[0]

        run1 = p_tab.add_run()
        run1.text = "단지에서 학교까지  "
        run1.font.size = Pt(9)
        run1.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        run2 = p_tab.add_run()
        run2.text = school_info.elementary_walk_distance
        run2.font.size = Pt(10)
        run2.font.bold = True
        run2.font.color.rgb = RGBColor(0x00, 0x7B, 0xD0)

    # 상세정보 테이블 (key-value)
    info_rows = []
    if school_info.elementary_address:
        info_rows.append(("주소", school_info.elementary_address))
    if school_info.elementary_phone:
        info_rows.append(("전화", school_info.elementary_phone))
    if school_info.elementary_founding_date:
        founding_text = school_info.elementary_type
        if school_info.elementary_founding_date:
            founding_text += f"  {school_info.elementary_founding_date}"
        info_rows.append(("설립", founding_text.strip()))
    if school_info.elementary_education_office:
        info_rows.append(("교육청", school_info.elementary_education_office))
    if school_info.elementary_coedu:
        info_rows.append(("구분", school_info.elementary_coedu))
    if school_info.elementary_homepage:
        info_rows.append(("홈페이지", school_info.elementary_homepage))

    if not info_rows:
        return

    top_y = 2.45
    row_height = 0.32
    label_width = Inches(0.7)
    value_width = Inches(3.1)
    left_x = Inches(5.8)

    for i, (label, value) in enumerate(info_rows):
        y = Inches(top_y + i * row_height)

        if i == 0:
            line = slide.shapes.add_shape(
                1, left_x, y - Pt(2), label_width + value_width, Pt(1)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
            line.line.fill.background()

        label_box = slide.shapes.add_textbox(left_x, y, label_width, Inches(row_height))
        lf = label_box.text_frame
        lf.word_wrap = True
        lp = lf.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(8)
        lp.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        value_box = slide.shapes.add_textbox(
            left_x + label_width, y, value_width, Inches(row_height)
        )
        vf = value_box.text_frame
        vf.word_wrap = True
        vp = vf.paragraphs[0]
        vp.text = value
        vp.font.size = Pt(8)
        vp.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        line_y = y + Inches(row_height) - Pt(2)
        line = slide.shapes.add_shape(
            1, left_x, line_y, label_width + value_width, Pt(1)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
        line.line.fill.background()


def add_middle_high_school_slide(
    prs: Presentation,
    complex_name: str,
    school_info: SchoolInfo,
    logo_path: Optional[str] = None,
):
    """학군지도(중·고등학교) 슬라이드 추가 (Slide D)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 공통 헤더
    add_slide_header(slide, "학군지도(중·고등학교)", "학군지도")

    # 학군 지도 이미지
    if (school_info.middle_high_map_path
            and os.path.exists(school_info.middle_high_map_path)):
        slide.shapes.add_picture(
            school_info.middle_high_map_path,
            Inches(0.5), Inches(1.2), Inches(9.0), Inches(4.3)
        )
    else:
        _add_placeholder_text(
            slide, Inches(0.5), Inches(1.2), Inches(9.0), Inches(4.3),
            "[중·고등학교 학군 지도]"
        )

    # 출처
    add_source_text(slide, "*아실")

    # 로고
    add_logo(slide, logo_path)

    return slide


def _add_placeholder_text(slide, left, top, width, height, text):
    """이미지 없을 때 텍스트 placeholder"""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER
