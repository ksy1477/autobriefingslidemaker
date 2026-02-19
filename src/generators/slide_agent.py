"""
중개사 소개 슬라이드 생성 — 프로필 모음 PPT에서 복사
"""
import os
import copy
from io import BytesIO
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.shapes.picture import Picture
from pptx.shapes.group import GroupShape


# 프로필 PPT 높이(5.62") → 출력 PPT 높이(6.25") 비율
_Y_SCALE = 6.25 / 5.62  # ≈ 1.112


def add_agent_slide_from_collection(
    prs: Presentation,
    agent_name: str,
    profile_pptx_path: str,
):
    """
    프로필 모음 PPT에서 이름 매칭하여 슬라이드 복사

    Args:
        prs: 출력 Presentation 객체
        agent_name: 중개사 이름 (매칭 키)
        profile_pptx_path: 프로필 모음 PPTX 경로
    """
    if not os.path.exists(profile_pptx_path):
        print(f"  [WARN] 프로필 모음 파일 없음: {profile_pptx_path}")
        _add_fallback_slide(prs, agent_name)
        return

    profile_prs = Presentation(profile_pptx_path)
    matched_slide = _find_matching_slide(profile_prs, agent_name)

    if matched_slide is None:
        print(f"  [WARN] 프로필 매칭 실패: '{agent_name}' 이름을 찾을 수 없습니다.")
        _add_fallback_slide(prs, agent_name)
        return

    # 빈 슬라이드 추가
    slide_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(slide_layout)

    # 매칭된 슬라이드의 모든 shape를 복사
    _copy_shapes(matched_slide, new_slide)

    print(f"  [INFO] 프로필 매칭 성공: '{agent_name}'")


def _find_matching_slide(profile_prs: Presentation, agent_name: str):
    """프로필 PPT에서 agent_name이 포함된 슬라이드 검색"""
    for slide in profile_prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                full_text = "".join(
                    run.text for para in shape.text_frame.paragraphs
                    for run in para.runs
                )
                if agent_name in full_text:
                    return slide
            if isinstance(shape, GroupShape):
                if _search_group_text(shape, agent_name):
                    return slide
    return None


def _search_group_text(group_shape, agent_name: str) -> bool:
    """GroupShape 내부에서 agent_name 검색"""
    for shape in group_shape.shapes:
        if shape.has_text_frame:
            full_text = "".join(
                run.text for para in shape.text_frame.paragraphs
                for run in para.runs
            )
            if agent_name in full_text:
                return True
        if isinstance(shape, GroupShape):
            if _search_group_text(shape, agent_name):
                return True
    return False


def _copy_shapes(src_slide, dst_slide):
    """원본 슬라이드의 모든 shape를 대상 슬라이드로 복사"""
    for shape in src_slide.shapes:
        if isinstance(shape, Picture):
            _copy_picture(shape, dst_slide)
        elif isinstance(shape, GroupShape):
            _copy_group_shape(shape, dst_slide)
        elif shape.has_text_frame:
            _copy_textbox(shape, dst_slide)
        else:
            _copy_shape_xml(shape, dst_slide)


def _scale_y(y_emu: int) -> int:
    """y좌표를 비율 조정 (5.62" → 6.25")"""
    return int(y_emu * _Y_SCALE)


def _copy_picture(shape: Picture, dst_slide):
    """이미지 shape 복사: blob 추출 → BytesIO → add_picture"""
    try:
        image = shape.image
        blob = image.blob
        stream = BytesIO(blob)

        left = shape.left
        top = _scale_y(shape.top)
        width = shape.width
        height = int(shape.height * _Y_SCALE)

        dst_slide.shapes.add_picture(stream, left, top, width, height)
    except Exception as e:
        print(f"    [WARN] 이미지 복사 실패: {e}")


def _copy_textbox(shape, dst_slide):
    """텍스트박스 복사: 위치/크기 + paragraph/run 폰트속성 전체 복사"""
    left = shape.left
    top = _scale_y(shape.top)
    width = shape.width
    height = int(shape.height * _Y_SCALE)

    new_txBox = dst_slide.shapes.add_textbox(left, top, width, height)
    new_tf = new_txBox.text_frame
    new_tf.word_wrap = shape.text_frame.word_wrap

    src_paragraphs = list(shape.text_frame.paragraphs)

    for i, src_para in enumerate(src_paragraphs):
        if i == 0:
            dst_para = new_tf.paragraphs[0]
        else:
            dst_para = new_tf.add_paragraph()

        dst_para.alignment = src_para.alignment
        if src_para.space_before is not None:
            dst_para.space_before = src_para.space_before
        if src_para.space_after is not None:
            dst_para.space_after = src_para.space_after

        for src_run in src_para.runs:
            dst_run = dst_para.add_run()
            dst_run.text = src_run.text
            _copy_font(src_run.font, dst_run.font)

    if shape.rotation:
        new_txBox.rotation = shape.rotation


def _copy_font(src_font, dst_font):
    """폰트 속성 복사"""
    if src_font.size is not None:
        dst_font.size = src_font.size
    if src_font.bold is not None:
        dst_font.bold = src_font.bold
    if src_font.italic is not None:
        dst_font.italic = src_font.italic
    if src_font.underline is not None:
        dst_font.underline = src_font.underline
    try:
        if src_font.color and src_font.color.rgb:
            dst_font.color.rgb = src_font.color.rgb
    except Exception:
        pass
    if src_font.name:
        dst_font.name = src_font.name


def _copy_group_shape(group: GroupShape, dst_slide):
    """GroupShape — 내부 shape를 개별적으로 복사 (그룹 해체)"""
    for shape in group.shapes:
        if isinstance(shape, Picture):
            _copy_picture(shape, dst_slide)
        elif shape.has_text_frame:
            _copy_textbox(shape, dst_slide)
        elif isinstance(shape, GroupShape):
            _copy_group_shape(shape, dst_slide)
        else:
            _copy_shape_xml(shape, dst_slide)


def _copy_shape_xml(shape, dst_slide):
    """XML 복사로 도형 복사 (y좌표 조정 포함)"""
    try:
        el = copy.deepcopy(shape._element)

        nsmap = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }
        xfrm = el.find('.//a:xfrm', nsmap)
        if xfrm is not None:
            off_el = xfrm.find('a:off', nsmap)
            if off_el is not None and off_el.get('y'):
                old_y = int(off_el.get('y'))
                off_el.set('y', str(_scale_y(old_y)))
            ext_el = xfrm.find('a:ext', nsmap)
            if ext_el is not None and ext_el.get('cy'):
                old_cy = int(ext_el.get('cy'))
                ext_el.set('cy', str(int(old_cy * _Y_SCALE)))

        dst_slide.shapes._spTree.append(el)
    except Exception as e:
        print(f"    [WARN] 도형 XML 복사 실패: {e}")


def _add_fallback_slide(prs: Presentation, agent_name: str):
    """매칭 실패 시 빈 슬라이드 + 경고 텍스트"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"중개파트너 프로필: {agent_name}"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "(프로필 모음에서 매칭되는 슬라이드를 찾을 수 없습니다)"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p2.alignment = PP_ALIGN.CENTER

    return slide
