"""
물건리스트 슬라이드 생성 — 레퍼런스 PPT 양식
"""
import os
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from src.models import PropertyDetail, ComplexData
from src.generators.slide_utils import add_slide_header, add_logo


def add_list_slide(
    prs: Presentation,
    complex_data_list: List[ComplexData],
    logo_path: Optional[str] = None,
):
    """
    물건리스트 슬라이드 추가 (레퍼런스 레이아웃)

    - Group 장식 마커 + 제목 (21pt bold)
    - 8열 테이블: 지역, 단지명, 동/호수, 평형(전용), 매매가격, 향/구조, 특이사항, 비고
    - 헤더: EFEFEF 배경, 검정 10pt bold
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 공통 헤더
    add_slide_header(slide, "물건리스트", "투어일정")

    # 전체 매물 목록 플랫화
    all_properties = []
    for cdata in complex_data_list:
        for prop in cdata.properties:
            all_properties.append((cdata, prop))

    # 8열 테이블
    cols = 8  # 지역, 단지명, 동/호수, 평형(전용), 매매가격, 향/구조, 특이사항, 비고
    rows = len(all_properties) + 1  # 헤더 + 데이터

    if rows < 2:
        rows = 2

    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(0.430), Inches(0.851),
        Inches(9.14), Inches(min(4.5, rows * 0.4))
    )
    table = table_shape.table

    # 열 너비 설정
    col_widths = [
        Inches(0.8),   # 지역
        Inches(2.0),   # 단지명
        Inches(1.0),   # 동/호수
        Inches(0.9),   # 평형(전용)
        Inches(1.0),   # 매매가격
        Inches(1.0),   # 향/구조
        Inches(1.7),   # 특이사항
        Inches(0.74),  # 비고
    ]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # 헤더: EFEFEF 배경, 검정 10pt bold
    headers = ["지역", "단지명", "동/호수", "평형(전용)", "매매가격", "향/구조", "특이사항", "비고"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xEF, 0xEF, 0xEF)

    # 데이터 행
    for row_idx, (cdata, prop) in enumerate(all_properties, start=1):
        address = cdata.complex_info.address
        addr_short = address.split()[:2] if address else [""]
        addr_text = " ".join(addr_short)

        dong_ho = prop.dong
        if hasattr(prop, 'ho') and getattr(prop, 'ho', None):
            dong_ho += f" {prop.ho}"

        # 단지명 형식: "한신 (95y, 1248^)" — built_year, total_units 활용
        ci = cdata.complex_info
        complex_display = ci.name
        year_suffix = f"{ci.built_year % 100}y" if ci.built_year else ""
        units_suffix = f"{ci.total_units}^" if ci.total_units else ""
        if year_suffix or units_suffix:
            parts = ", ".join(filter(None, [year_suffix, units_suffix]))
            complex_display = f"{ci.name} ({parts})"

        area_text = prop.area_pyeong or ""
        direction_structure = " ".join(
            filter(None, [prop.direction, prop.structure])
        )

        row_data = [
            addr_text,
            complex_display,
            dong_ho,
            area_text,
            prop.price,
            direction_structure,
            prop.memo or "",
            "",  # 비고
        ]

        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                paragraph.alignment = PP_ALIGN.CENTER

            # 짝수 행 배경색
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    # 로고
    add_logo(slide, logo_path)

    return slide
