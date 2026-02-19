"""
표지 슬라이드 생성 — 레퍼런스 PPT 양식
"""
import os
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def add_cover_slide(
    prs: Presentation,
    customer_name: str,
    complex_names: List[str],
    logo_path: Optional[str] = None,
    background_path: Optional[str] = None,
):
    """
    표지 슬라이드 추가 (레퍼런스 레이아웃)

    레이아웃:
    - 좌측: 고객명+브리핑제목 (우측정렬, 40pt bold) at (-0.048", 1.422")
    - 좌측: 임장브리핑 텍스트 (18pt) at (0.105", 4.095")
    - 좌하단: 로고 이미지 at (1.357", 4.770", 3.033"×1.480")
    - 우측: cover_image.png (5.952", 0", 4.048"×6.25")
    """
    slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
    slide = prs.slides.add_slide(slide_layout)

    # 우측 커버 이미지
    cover_image_path = background_path or os.path.join("assets", "cover_image.png")
    if os.path.exists(cover_image_path):
        slide.shapes.add_picture(
            cover_image_path,
            Inches(5.952), Inches(0),
            Inches(4.048), Inches(6.25),
        )

    # 좌측 고객명 + 브리핑 제목 (우측정렬, 40pt bold)
    complex_text = " / ".join(complex_names) if complex_names else ""

    txBox = slide.shapes.add_textbox(
        Inches(-0.048), Inches(1.422),
        Inches(5.8), Inches(2.5),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    # 고객명
    p1 = tf.paragraphs[0]
    p1.text = f"{customer_name}님"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(0x22, 0x1E, 0x1F)
    p1.alignment = PP_ALIGN.RIGHT

    # 단지명들
    p2 = tf.add_paragraph()
    p2.text = complex_text
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0x22, 0x1E, 0x1F)
    p2.alignment = PP_ALIGN.RIGHT

    # 임장브리핑 텍스트 (18pt)
    txBox2 = slide.shapes.add_textbox(
        Inches(0.105), Inches(4.095),
        Inches(5.5), Inches(0.5),
    )
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = "임장브리핑"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(0x22, 0x1E, 0x1F)
    p3.alignment = PP_ALIGN.RIGHT

    # 좌하단 로고 이미지
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(
            logo_path,
            Inches(1.357), Inches(4.770),
            Inches(3.033), Inches(1.480),
        )

    return slide
